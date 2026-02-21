import json
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt

TEMPLATE_PPTX = "template/report_template.pptx"
OUTPUT_DIR = "output"

# Update this to the number of milestone rows you created in the template slide
MAX_MILESTONES = 5

# --- helpers ---
def find_shape_by_name(slide, shape_name: str):
    for shp in slide.shapes:
        if shp.name == shape_name:
            return shp
    return None

def set_text(slide, shape_name: str, value: str):
    shp = find_shape_by_name(slide, shape_name)
    if shp is None:
        raise KeyError(f"Shape not found in template slide: '{shape_name}'")
    if not shp.has_text_frame:
        raise TypeError(f"Shape '{shape_name}' does not support text (no text_frame).")

    tf = shp.text_frame
    text = "" if value is None else str(value)

    # If there's at least one run, just replace text in the first run (preserves style)
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        # Remove any extra runs (sometimes template has multiple)
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    else:
        # fallback: create one run, but still avoid tf.clear() when possible
        if not tf.paragraphs:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text

def set_lines(slide, shape_name: str, lines):
    """
    Writes a list of lines into a textbox, each line in a new paragraph.
    """
    shp = find_shape_by_name(slide, shape_name)
    if shp is None:
        raise KeyError(f"Shape not found in template slide: '{shape_name}'")
    if not shp.has_text_frame:
        raise TypeError(f"Shape '{shape_name}' does not support text (no text_frame).")

    tf = shp.text_frame
    tf.clear()

    if not lines:
        tf.text = ""
        return

    # first line
    tf.text = str(lines[0])
    # remaining lines
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = str(line)

def shrink_font_if_long(slide, shape_name: str, max_chars=220, start_pt=18, min_pt=10):
    """
    Very basic font shrink heuristic to reduce overflow risk.
    """
    shp = find_shape_by_name(slide, shape_name)
    if shp is None or not shp.has_text_frame:
        return

    text = shp.text_frame.text or ""
    n = len(text)

    if n <= max_chars:
        size = start_pt
    else:
        # linearly shrink down to min_pt
        # cap so it doesn't go below min_pt
        extra = min(n - max_chars, 400)
        size = max(min_pt, start_pt - int(extra / 80) - 1)

    for p in shp.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)

def duplicate_slide(prs: Presentation, slide_index: int):
    """
    Duplicates a slide by deep-copying its XML tree (preserves layout/formatting).
    """
    source = prs.slides[slide_index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Remove default placeholders in new slide
    for shp in list(new_slide.shapes):
        el = shp._element
        el.getparent().remove(el)

    # Copy shapes from source
    for shp in source.shapes:
        new_el = deepcopy(shp._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide

def fill_project_slide(slide, report_date: str, project: dict):
    # Main fields
    set_text(slide, "Project_title", project.get("title", ""))
    set_text(slide, "project_scope_details", project.get("scope", ""))
    set_text(slide, "project_user", project.get("user", ""))
    set_text(slide, "project_area", project.get("area", ""))
    set_text(slide, "report_date", report_date)

    # Status fields
    set_text(slide, "progress_percentage_", f'{project.get("progress_percent", 0)}%')
    set_text(slide, "days_total_", str(project.get("days_total", "")))

    # Observations / Risks (one per line)
    set_lines(slide, "observations_", project.get("observations", []))
    set_lines(slide, "risk_", project.get("risks", []))

    # Optional: shrink long text to reduce overflow
    shrink_font_if_long(slide, "project_scope_details", max_chars=250, start_pt=14, min_pt=12)
    shrink_font_if_long(slide, "observations_", max_chars=220, start_pt=11, min_pt=9)
    shrink_font_if_long(slide, "risk_", max_chars=220, start_pt=11, min_pt=10)

    # Milestones mapping (your 4-field row)
    milestones = project.get("milestones", [])

    for i in range(1, MAX_MILESTONES + 1):
        ms = milestones[i - 1] if i - 1 < len(milestones) else {
            "milestone": "",
            "date_log": "",
            "status_log": "",
            "status": ""
        }

        set_text(slide, f"milestone_{i}", ms.get("milestone", ""))
        set_text(slide, f"date_log{i}", ms.get("date_log", ""))
        set_text(slide, f"status_log{i}", ms.get("status_log", ""))
        set_text(slide, f"status{i}", ms.get("status", ""))

def generate(json_path: str):
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    report_date = data.get("report_date", "")
    projects = data.get("projects", [])

    if not projects:
        raise ValueError("No projects found in JSON. Expected data['projects'] to be a non-empty list.")

    prs = Presentation(TEMPLATE_PPTX)

    # We assume slide 0 is the template slide
    template_index = 0

    # Ensure we have only 1 template slide at start (optional)
    # If your template file already has multiple slides, you can remove this block
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    # Fill first slide with first project
    fill_project_slide(prs.slides[0], report_date, projects[0])

    # Create additional slides for remaining projects
    for project in projects[1:]:
        slide = duplicate_slide(prs, template_index)
        fill_project_slide(slide, report_date, project)

    out_path = os.path.join(OUTPUT_DIR, f"Weekly_Report_{report_date}.pptx")
    prs.save(out_path)
    print(f"✅ Generated: {out_path}")

if __name__ == "__main__":
    # Change this to your file name if needed
    generate("data/2026-02-20.json")